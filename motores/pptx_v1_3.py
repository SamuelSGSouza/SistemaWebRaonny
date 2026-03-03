from pptx import Presentation
from pptx.oxml.xmlchemy import OxmlElement
import copy, string
from pptx.table import Table
from pptx.oxml.table import CT_TableRow
from datetime import datetime
from io import BytesIO
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls
import subprocess, secrets, os
from pptx.util import Emu
from pptx.enum.text import MSO_AUTO_SIZE

def remove_table_borders(table):
    """Remove as bordas de todas as células de uma tabela."""
    for row in table.rows:
        for cell in row.cells:
            # Acessa as propriedades da célula no XML
            tcPr = cell._tc.get_or_add_tcPr()
            
            # Lista de todas as bordas possíveis
            borders = ['lnL', 'lnR', 'lnT', 'lnB', 'lnTlToBr', 'lnBlToTr']
            
            for border in borders:
                # Cria um elemento de "no fill" (sem preenchimento) para a borda
                ln = OxmlElement(f'a:{border}')
                noFill = OxmlElement('a:noFill')
                ln.append(noFill)
                tcPr.append(ln)

def clonar_objeto_xml(elemento:CT_TableRow) -> CT_TableRow:
    """Cria uma cópia profunda de um elemento XML."""
    return copy.deepcopy(elemento)

def substituicao_parte_1(prs,to_replace, replace):
    """Varre todos os slides e substitui o placeholder pelo nome da empresa."""
    for slide in prs.slides:
        for shape in slide.shapes:
            # 1. Substituir em formas com texto (caixas de texto, retângulos, etc)
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if to_replace in run.text:
                            run.text = run.text.replace(to_replace, replace)

            # 2. Substituir dentro de tabelas
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if to_replace in run.text:
                                    run.text = run.text.replace(to_replace, replace)

def replace_text_in_paragraph(paragraph, subs):
    inline = paragraph.runs
    for search_str, replace_str in subs.items():
        if search_str in paragraph.text:
            # Garante que o replace_str seja string e remove o caractere problemático
            texto_limpo = str(replace_str).replace('\x0b', '') 
            
            for i in range(len(inline)):
                run = inline[i]
                if search_str in run.text:
                    run.text = run.text.replace(search_str, texto_limpo)
                else:
                    # No caso de reconstrução do parágrafo:
                    full_text = paragraph.text.replace(search_str, texto_limpo)
                    
                    # Salva a formatação do primeiro run
                    if len(paragraph.runs) > 0:
                        primeiro_run = paragraph.runs[0]
                        font = primeiro_run.font
                        
                        # Limpa todos os runs e cria um novo com o texto completo
                        # Isso mantém a formatação da primeira palavra para o parágrafo todo
                        p_text = full_text
                        for r in paragraph.runs:
                            r.text = ""
                        
                        novo_run = paragraph.add_run()
                        novo_run.text = p_text
                        
                        # Copia propriedades básicas de fonte
                        try:
                            novo_run.font.name = font.name
                            novo_run.font.size = font.size
                            novo_run.font.bold = font.bold
                            novo_run.font.italic = font.italic
                            novo_run.font.color.rgb = font.color.rgb
                        except:
                            pass # Algumas propriedades podem não estar definidas
                    break

def substituicao_parte_2(shape, subs):
    """Recursividade para tratar grupos, tabelas e formas simples."""
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            replace_text_in_paragraph(paragraph, subs)

    elif shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                for paragraph in cell.text_frame.paragraphs:
                    replace_text_in_paragraph(paragraph, subs)

    elif shape.shape_type == 6:  # GroupShape
        for sub_shape in shape.shapes:
            substituicao_parte_2(sub_shape, subs)

def aplicar_cor_fundo(cell, cor_hex):
    """Aplica uma cor de fundo (HEX) a uma célula."""
    # Remove preenchimento existente se houver
    tcPr = cell._tc.get_or_add_tcPr()
    for child in tcPr:
        if child.tag.endswith('solidFill'):
            tcPr.remove(child)
            
    # Cria o elemento de preenchimento sólido
    fill_xml = f'<a:solidFill {nsdecls("a")}><a:srgbClr val="{cor_hex}"/></a:solidFill>'
    tcPr.append(parse_xml(fill_xml))

def aplicar_bordas(cell, cor_hex="000000", largura="12700"):
    """Aplica bordas simples nos quatro lados da célula."""
    tcPr = cell._tc.get_or_add_tcPr()
    
    # XML para as 4 bordas (lnL, lnR, lnT, lnB)
    # largura 12700 é aproximadamente 1pt
    bordas_xml = f"""
    <a:tcPr {nsdecls("a")}>
        <a:lnL w="{largura}"><a:solidFill><a:srgbClr val="{cor_hex}"/></a:solidFill></a:lnL>
        <a:lnR w="{largura}"><a:solidFill><a:srgbClr val="{cor_hex}"/></a:solidFill></a:lnR>
        <a:lnT w="{largura}"><a:solidFill><a:srgbClr val="{cor_hex}"/></a:solidFill></a:lnT>
        <a:lnB w="{largura}"><a:solidFill><a:srgbClr val="{cor_hex}"/></a:solidFill></a:lnB>
    </a:tcPr>
    """
    # Mescla as propriedades de borda
    for border_element in parse_xml(bordas_xml):
        tcPr.append(border_element)

def estilizar_tabela_toda(tabela, cor_par="FFFFFF", cor_impar="D9D9D9", cor_borda="767676"):
    """
    Percorre todas as linhas da tabela e aplica a estilização intercalada e bordas.
    Ignora a primeira linha (geralmente o cabeçalho).
    """
    for idx, row in enumerate(tabela.rows):
        # Se quiser pular o cabeçalho, use: if idx == 0: continue
        # Caso queira estilizar inclusive o cabeçalho, remova a condição abaixo
        if idx == 0: 
            # Opcional: Aplicar apenas bordas no cabeçalho
            for cell in row.cells:
                aplicar_bordas(cell, cor_hex=cor_borda)
            continue

        # Define a cor baseada na paridade da linha
        cor = cor_par if idx % 2 == 0 else cor_impar
        
        for cell in row.cells:
            aplicar_cor_fundo(cell, cor)
            aplicar_bordas(cell, cor_hex=cor_borda)

def adicionar_linha_com_estilo(tabela: Table, dados: list):
    """Adiciona a linha e o texto. A estilização será feita depois pela estilizar_tabela_toda."""
    if not dados:
        return
    
    tbl = tabela._tbl
    ultima_linha_xml = tbl.tr_lst[-1]
    nova_linha_xml = clonar_objeto_xml(ultima_linha_xml)
    
    for tc in nova_linha_xml.tc_lst:
        for p in tc.xpath('.//a:p'):
            runs = p.xpath('.//a:r')
            if runs:
                runs[0].text = "" 
                for r_extra in runs[1:]:
                    p.remove(r_extra)

    tbl.append(nova_linha_xml)
    nova_linha = tabela.rows[len(tabela.rows)-1]

    for i, texto in enumerate(dados):
        if i < len(nova_linha.cells):
            cell = nova_linha.cells[i]
            paragraph = cell.text_frame.paragraphs[0]
            if paragraph.runs:
                paragraph.runs[0].text = str(texto)
            else:
                run = paragraph.add_run()
                run.text = str(texto)
    

def ajustar_shapes_abaixo(slide, tabela_shape, delta_altura):
    base_inferior = tabela_shape.top + tabela_shape.height

    for shape in slide.shapes:
        if shape is tabela_shape:
            continue

        # Se o shape começa abaixo da tabela, move para baixo
        if shape.top >= base_inferior - delta_altura:
            shape.top = Emu(shape.top + delta_altura)

def remover_coluna(tabela, col_idx):
    """Remove uma coluna da tabela baseada no índice."""
    tbl = tabela._tbl
    
    # 1. Remover a definição da coluna (largura) no XML
    gridCol_lst = tbl.tblGrid.gridCol_lst
    if col_idx < len(gridCol_lst):
        gridCol = gridCol_lst[col_idx]
        tbl.tblGrid.remove(gridCol)
    
    # 2. Remover a célula (tc) de cada linha (tr)
    for tr in tbl.tr_lst:
        tc_lst = tr.tc_lst
        if col_idx < len(tc_lst):
            tc = tc_lst[col_idx]
            tr.remove(tc)

def remover_linha(tabela, row_idx):
    """Remove uma linha da tabela baseada no índice."""
    tbl = tabela._tbl
    tr_lst = tbl.tr_lst # Lista de todas as linhas (elementos <w:tr>)
    
    if row_idx < len(tr_lst):
        tr = tr_lst[row_idx]
        tbl.remove(tr)
    else:
        print(f"Erro: O índice {row_idx} está fora do alcance da tabela.")

def processar_apresentacao(caminho_entrada, subs:dict, servicos:list, equipamentos:list, adicionais:list, pagina_servicos:int=0, pagina_adicionais:int=0,pagina_equipamentos:int=0):
    try:
        prs = Presentation(caminho_entrada)
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    # Isso impede que o PowerPoint tente "encolher" o texto automaticamente
                    shape.text_frame.word_wrap = True
                    shape.text_frame.auto_size = MSO_AUTO_SIZE.NONE

        for k,v in subs.items():
            substituicao_parte_1(prs, k, v)

        # for section in [prs.slides, prs.slide_masters, prs.slide_layouts]:
        #     for item in section:
        #         for shape in item.shapes:
        #             substituicao_parte_2(shape, subs)

        contador_paginas = 0
        for slide in prs.slides:
            contador_paginas += 1
            conj_dados = []
            
            if contador_paginas == pagina_equipamentos:
                conj_dados = equipamentos
            elif contador_paginas == pagina_servicos:
                conj_dados = servicos
            elif contador_paginas == pagina_adicionais:
                conj_dados = adicionais
                
            for shape in slide.shapes:
                
                if shape.has_table:
                    linhas_antes = len(shape.table.rows)

                    for dados in conj_dados:
                        adicionar_linha_com_estilo(shape.table, dados)

                    if conj_dados != []:
                        remover_linha(shape.table, 1)

                    linhas_depois = len(shape.table.rows)
                    linhas_adicionadas = linhas_depois - linhas_antes

                    # if linhas_adicionadas > 0:
                    #     # soma a altura real das linhas recém-criadas
                    #     delta = 0
                    #     for i in range(linhas_antes, linhas_depois):
                    #         delta += shape.table.rows[i].height

                    #     # ajusta o tamanho do shape exatamente para o necessário
                    #     shape.height = Emu(shape.height + delta)

                    #     ajustar_shapes_abaixo(slide, shape, delta)
                    if contador_paginas == 1:
                        remove_table_borders(shape.table)
                    else:
                        estilizar_tabela_toda(shape.table)
        return prs
    except Exception as e:
        print(f"Erro inesperado: {e}")

def gerador_modelo_2(
        nome_empresa,
        cnpj_empresa,
        telefone_empresa, 
        cidade_empresa,
        uf_empresa, 
        nome_responsavel, 
        email_responsavel, 
        sexo_responsavel, 
        numero_proposta, 
        valor_dolar,
        servicos,
        equipamentos,
        adicionais,
        pagina_equipamentos,
        pagina_servicos,
        pagina_adicionais,
        modelo_path,
        tempo_contrato,
        nome_usuario_responsavel,
        telefone_usuario_responsavel,
        email_usuario_responsavel,
        cargo_responsavel,
        observacao_adicional,
        observacao_servico,
        observacao_equipamento,
        formato_download="pptx",
):
    subs = {
        "{{nome_empresa}}": nome_empresa,
        "{{cnpj_empresa}}": cnpj_empresa,
        "{{telefone_empresa}}": telefone_empresa,
        "{{cidade_empresa}}": cidade_empresa,
        "{{uf_empresa}}": uf_empresa,
        "{{nome_responsavel}}": nome_responsavel,
        "{{email_responsavel}}": email_responsavel,
        "{{cargo_usuario_responsavel}}": cargo_responsavel,
        "{{sexo_responsavel}}": sexo_responsavel,
        "{{data}}": datetime.today().strftime("%d/%m/%Y"),
        "{{valor_dolar}}": valor_dolar,
        "{{numero_proposta}}": numero_proposta,
        "{{tempo_de_contrato}}": tempo_contrato,
        "{{nome_usuario_responsavel}}": nome_usuario_responsavel,
        "{{telefone_usuario_responsavel}}": telefone_usuario_responsavel,
        "{{email_usuario_responsavel}}": email_usuario_responsavel,
        "{{observacao_tabela_variavel}}": observacao_adicional,
        "{{observacao_tabela_equipamento}}": observacao_equipamento,
        "{{observacao_tabela_servico}}": observacao_servico,

    }
    doc = processar_apresentacao(modelo_path,  subs, servicos=servicos, equipamentos=equipamentos,adicionais=adicionais, pagina_adicionais=pagina_adicionais, pagina_equipamentos=pagina_equipamentos, pagina_servicos=pagina_servicos)
    if formato_download == "pdf":
        token = secrets.token_hex(10)
        pptx_path = f"/tmp/{token}.pptx"
        pdf_path = f"/tmp/{token}.pdf"

        try:
            doc.save(pptx_path)
            converter_para_pdf(pptx_path, "/tmp")

            with open(pdf_path, "rb") as f:
                pdf_bytes = f.read()

            return pdf_bytes

        finally:
            # Remove os arquivos se existirem
            for path in (pptx_path, pdf_path):
                try:
                    if os.path.exists(path):
                        os.remove(path)
                except Exception:
                    # Em ambiente Lambda, por exemplo, falhas aqui não devem quebrar a resposta
                    pass
    
    buffer = BytesIO()
    doc.save(buffer)
    buffer.seek(0)

    bytes_doc = buffer.getvalue()
    return bytes_doc


def converter_para_pdf(caminho_pptx, pasta_saida):
    subprocess.run([
        "libreoffice",
        "--headless",
        "--convert-to", "pdf",
        "--outdir", pasta_saida,
        caminho_pptx
    ], check=True)

# prs = gerador_modelo_2()
# pptx_gerado = "/tmp/apresentacao.pptx"
# prs.save(pptx_gerado)

